# -*- coding: utf-8 -*-
"""Genera PRESENTAZIONE_v2.pptx: la cronologia del progetto in 19 slide (16:9)."""
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
FIG = ROOT / "figures"
OUT = ROOT / "PRESENTAZIONE_v2.pptx"

INK = RGBColor(0x1C, 0x21, 0x28)
DARK = RGBColor(0x16, 0x32, 0x4A)
BLUE = RGBColor(0x2D, 0x6A, 0x8F)
GREEN = RGBColor(0x1F, 0x9D, 0x55)
RED = RGBColor(0xC2, 0x55, 0x3A)
GREY = RGBColor(0x6B, 0x74, 0x7E)
LIGHT = RGBColor(0xF2, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = 13.333, 7.5
FONT = "Segoe UI"
MONO = "Cascadia Mono"

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def style(run, size, color=INK, bold=False, italic=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def para(tf, runs, first=False, align=PP_ALIGN.LEFT, space_after=6, space_before=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        style(r, **opts)
    return p


def bullets(tf, items, size=15, first=True, gap=10, glyph_color=BLUE):
    for i, item in enumerate(items):
        runs = [("▪  ", dict(size=size, color=glyph_color, bold=True))]
        if isinstance(item, str):
            runs.append((item, dict(size=size, color=INK)))
        else:  # lista di run (testo, opts)
            runs.extend(item)
        para(tf, runs, first=(first and i == 0), space_after=gap)


def rect(slide, x, y, w, h, fill, rounded=True, line=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def pic_fit(slide, path, x, y, maxw, maxh):
    w, h = Image.open(path).size
    s = min(maxw / w, maxh / h)
    pw, ph = w * s, h * s
    return slide.shapes.add_picture(
        str(path), Inches(x + (maxw - pw) / 2), Inches(y + (maxh - ph) / 2),
        Inches(pw), Inches(ph))


def content_slide(kicker, title, accent=BLUE):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, SW, 0.12, accent, rounded=False)
    tf = textbox(slide, 0.6, 0.34, 12.1, 0.4)
    para(tf, [(kicker.upper(), dict(size=12, color=GREY, bold=True))], first=True)
    tf = textbox(slide, 0.6, 0.66, 12.1, 0.9)
    para(tf, [(title, dict(size=29, color=DARK, bold=True))], first=True)
    return slide


def takeaway(slide, runs, y=6.65):
    bar = rect(slide, 0.6, y, 12.13, 0.62, LIGHT)
    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    para(tf, runs, first=True, space_after=0)


def card(slide, x, y, w, h, number, label, color=BLUE, num_size=30):
    c = rect(slide, x, y, w, h, LIGHT)
    tf = c.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    para(tf, [(number, dict(size=num_size, color=color, bold=True))],
         first=True, align=PP_ALIGN.CENTER, space_after=2)
    para(tf, [(label, dict(size=11.5, color=GREY))],
         align=PP_ALIGN.CENTER, space_after=0)


# ---------------------------------------------------------------- 1 · titolo
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, 0.18, BLUE, rounded=False)
rect(s, 0, SH - 0.18, SW, 0.18, GREEN, rounded=False)
tf = textbox(s, 1.0, 1.55, 11.33, 1.3)
para(tf, [("Da OCGAN a PatchCore", dict(size=46, color=DARK, bold=True))],
     first=True, align=PP_ALIGN.CENTER)
tf = textbox(s, 1.6, 2.75, 10.13, 1.0)
para(tf, [("Anomaly detection one-class: da un paper CVPR 2019 "
           "a un sistema industriale su MVTec AD",
           dict(size=18, color=GREY))], first=True, align=PP_ALIGN.CENTER)
tf = textbox(s, 1.6, 3.62, 10.13, 0.5)
para(tf, [("Alessandro Pata  ·  giugno 2026", dict(size=14, color=INK))],
     first=True, align=PP_ALIGN.CENTER)
cw, gap = 3.4, 0.45
x0 = (SW - 3 * cw - 2 * gap) / 2
card(s, x0, 4.7, cw, 1.5, "11.738", "run di ablation su disco", BLUE)
card(s, x0 + cw + gap, 4.7, cw, 1.5, "15", "categorie MVTec AD", BLUE)
card(s, x0 + 2 * (cw + gap), 4.7, cw, 1.5, "0.9846", "macro AUROC finale", GREEN)

# ---------------------------------------------------------------- 2 · paper
s = content_slide("2019 · Il punto di partenza", "OCGAN — One-Class Novelty Detection (CVPR 2019)")
tf = textbox(s, 0.6, 1.75, 7.4, 4.0)
bullets(tf, [
    [("L'idea: ", dict(size=15, color=INK, bold=True)),
     ("non basta ricostruire bene la classe normale — l'intero spazio latente "
      "deve rappresentare solo la normalità", dict(size=15, color=INK))],
    [("Informative-negative mining: ", dict(size=15, color=INK, bold=True)),
     ("ricerca adversariale (gradient ascent nel latent) delle zone che generano "
      "immagini fuori classe — e lì si rieduca il generatore", dict(size=15, color=INK))],
    [("4 componenti: ", dict(size=15, color=INK, bold=True)),
     ("denoising autoencoder + latent discriminator + visual discriminator + "
      "classificatore che guida il mining", dict(size=15, color=INK))],
    [("Score finale: ", dict(size=15, color=INK, bold=True)),
     ("solo errore di ricostruzione (MSE), loss = 10·MSE + termini adversariali",
      dict(size=15, color=INK))],
])
p = rect(s, 8.35, 1.75, 4.35, 3.6, LIGHT)
tf = p.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.15)
para(tf, [("I numeri del paper", dict(size=13, color=GREY, bold=True))], first=True, space_after=8)
para(tf, [("MNIST    ", dict(size=16, color=INK)), ("0.9750", dict(size=20, color=BLUE, bold=True))], space_after=4)
para(tf, [("CIFAR-10  ", dict(size=16, color=INK)), ("0.6566", dict(size=20, color=RED, bold=True))], space_after=10)
para(tf, [("Ablation: l'autoencoder da solo fa già 0.957; "
           "tutto l'apparato adversariale aggiunge +0.018",
           dict(size=12.5, color=GREY, italic=True))], space_after=0)
takeaway(s, [("Limiti: ", dict(size=13.5, color=DARK, bold=True)),
             ("immagini piccole mono-concetto, score solo pixel, risultati misurati una volta sola.",
              dict(size=13.5, color=INK))], y=5.7)

# ---------------------------------------------------------------- 3 · obiettivo
s = content_slide("L'obiettivo", "Dal benchmark accademico al problema industriale")
tf = textbox(s, 0.6, 1.85, 12.1, 4.6)
bullets(tf, [
    [("MVTec AD: ", dict(size=16, color=INK, bold=True)),
     ("15 categorie (10 oggetti + 5 texture), training solo su pezzi integri, "
      "difetti reali al test (graffi, contaminazioni, parti mancanti…)", dict(size=16, color=INK))],
    [("Protocollo a 4 split: ", dict(size=16, color=INK, bold=True)),
     ("train_normal / val_normal / val_mixed / test blind — il test non tocca mai il tuning",
      dict(size=16, color=INK))],
    [("Multi-seed obbligatorio: ", dict(size=16, color=INK, bold=True)),
     ("3–5 seed, media ± std (su toothbrush la std del GAN è ±0.17: "
      "un seed solo racconta qualunque storia)", dict(size=16, color=INK))],
    [("Metrica guida: ", dict(size=16, color=INK, bold=True)),
     ("macro AUROC sulle 15 categorie, più AUPRC, F1 e FPR@95TPR", dict(size=16, color=INK))],
], gap=14)

# ---------------------------------------------------------------- 4 · binario B
s = content_slide("La strategia", "Binario B: modernizzare OCGAN, non difenderlo")
p = rect(s, 0.6, 1.8, 5.9, 3.95, LIGHT)
tf = p.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.16); tf.margin_right = Inches(0.18)
para(tf, [("DA OCGAN TENIAMO", dict(size=13, color=RED, bold=True))], first=True, space_after=10)
bullets(tf, [
    "La ricostruzione come segnale base",
    "Il controllo dello spazio latente (diventa compactness Deep-SVDD)",
    "Il mining dei negativi informativi (3 step PGD-style)",
], size=15, first=False, gap=10, glyph_color=RED)
p = rect(s, 6.85, 1.8, 5.9, 3.95, LIGHT)
tf = p.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.16); tf.margin_right = Inches(0.18)
para(tf, [("INNESTIAMO DAL MODERNO", dict(size=13, color=BLUE, bold=True))], first=True, space_after=10)
bullets(tf, [
    "Backbone ImageNet congelato (ResNet50)",
    "Loss percettive: L1 + MS-SSIM + perceptual (via MSE)",
    "Anomalie sintetiche: CutPaste, Perlin, rumore in feature space",
    "Memory bank + teacher–student come segnali ausiliari",
    "7 score fusi con regressione logistica",
], size=15, first=False, gap=8, glyph_color=BLUE)
takeaway(s, [("Riferimenti moderni: ", dict(size=13.5, color=DARK, bold=True)),
             ("PaDiM, PatchCore, DRAEM, RD4AD, EfficientAD, SimpleNet.",
              dict(size=13.5, color=INK))], y=6.1)

# ---------------------------------------------------------------- 5 · paperspace
s = content_slide("Infrastruttura", "Paperspace Gradient: la GPU in affitto")
tf = textbox(s, 0.6, 1.8, 7.5, 4.5)
bullets(tf, [
    [("Notebook cloud con GPU dedicata a costo orario", dict(size=15.5, color=INK, bold=True)),
     (" — ambiente CUDA pronto, zero amministrazione", dict(size=15.5, color=INK))],
    [("Storage persistente: ", dict(size=15.5, color=INK, bold=True)),
     ("dataset, checkpoint e 11.738 directory di run sopravvivono alle istanze effimere",
      dict(size=15.5, color=INK))],
    [("Fire-and-forget: ", dict(size=15.5, color=INK, bold=True)),
     ("lanci la griglia la sera, l'istanza macina da sola; paghi solo le ore di training",
      dict(size=15.5, color=INK))],
    [("Il contro, vissuto: ", dict(size=15.5, color=RED, bold=True)),
     ("sync del codice fuori da git → drift (+1.235 righe mai committate) → "
      "calibrazioni d'archivio non riproducibili", dict(size=15.5, color=INK))],
], gap=13)
p = rect(s, 8.45, 1.8, 4.28, 4.3, LIGHT)
tf = p.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.16)
para(tf, [("LE MACCHINE", dict(size=13, color=GREY, bold=True))], first=True, space_after=10)
para(tf, [("RTX A4000 16 GB", dict(size=15, color=DARK, bold=True))], space_after=2)
para(tf, [("griglia 11.7k run + final + production", dict(size=12, color=GREY))], space_after=10)
para(tf, [("Quadro RTX 5000 16 GB", dict(size=15, color=DARK, bold=True))], space_after=2)
para(tf, [("retrain ottimizzato optv2", dict(size=12, color=GREY))], space_after=10)
para(tf, [("Quadro T1000 (locale)", dict(size=15, color=DARK, bold=True))], space_after=2)
para(tf, [("inferenza live della webapp (fp32)", dict(size=12, color=GREY))], space_after=0)

# ---------------------------------------------------------------- 6 · architettura GAN
s = content_slide("Fase 1 · GAN modernizzato", "L'architettura: un OCGAN con organi nuovi")
p = rect(s, 0.6, 1.8, 6.6, 4.55, LIGHT)
tf = p.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.18)
flow = [
    "immagine 256×256",
    "ResNet50 pre-addestrata — CONGELATA",
    "latent compatto z ∈ R¹²⁸  (compactness 0.1)",
    "decoder: 5 × [upsample ×2 → conv 3×3 → BN → ReLU]",
    "ricostruzione 256×256",
    "ri-encoding nello stesso backbone",
]
for i, step in enumerate(flow):
    para(tf, [(step, dict(size=14.5, color=DARK if i in (1, 5) else INK,
                          bold=i in (1, 5), font=MONO))],
         first=(i == 0), align=PP_ALIGN.CENTER, space_after=1)
    if i < len(flow) - 1:
        para(tf, [("↓", dict(size=12, color=BLUE, bold=True))],
             align=PP_ALIGN.CENTER, space_after=1)
p = rect(s, 7.55, 1.8, 5.18, 4.55, LIGHT)
tf = p.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.18)
para(tf, [("7 SEGNALI DI ANOMALIA", dict(size=13, color=BLUE, bold=True))], first=True, space_after=8)
for name in ("ricostruzione (L1 + MS-SSIM)", "perceptual", "feature", "latente",
             "memory bank (kNN su patch)", "discriminativo (anomalie sintetiche)",
             "teacher–student (discrepanza)"):
    para(tf, [("▪  ", dict(size=13.5, color=BLUE, bold=True)),
              (name, dict(size=13.5, color=INK))], space_after=4)
para(tf, [("fusi da una regressione logistica su val_mixed, "
           "normalizzazione robusta MAD su val_normal",
           dict(size=12.5, color=GREY, italic=True))], space_before=8, space_after=0)

# ---------------------------------------------------------------- 7 · ablation
s = content_slide("Fase 1 · Ablation", "3.300 run: che cosa serve davvero?")
pic_fit(s, FIG / "fig2_fattori_ablation.png", 0.6, 1.65, 12.13, 4.75)
takeaway(s, [("La fusione appresa è IL fattore: +0.22 AUROC. ", dict(size=13.5, color=DARK, bold=True)),
             ("Lo score one-class — l'erede più diretto di OCGAN — toglie 0.06 e viene spento: "
              "i dati battono l'affetto per le idee.", dict(size=13.5, color=INK))])

# ---------------------------------------------------------------- 8 · reset onesto
s = content_slide("Fase 1 · Il bug più importante", "Il reset dell'onestà", accent=RED)
tf = textbox(s, 0.6, 1.85, 7.9, 4.4)
bullets(tf, [
    [("Tre flag di config “morti”: ", dict(size=15.5, color=INK, bold=True)),
     ("use_skip_connections, unfreeze_from, scoring_topk — letti dal config, "
      "ignorati dal builder del modello", dict(size=15.5, color=INK))],
    "Esperimenti che “confrontavano architetture diverse” stavano rieseguendo la stessa rete",
    "Scoperto ispezionando i pesi dei checkpoint, non i log",
    [("Fix → re-test completo → azzeramento dei numeri storici",
      dict(size=15.5, color=INK, bold=True))],
    [("Lezione: ", dict(size=15.5, color=RED, bold=True)),
     ("un config accettato non è un config applicato — ora un test verifica che ogni "
      "chiave abbia un effetto misurabile", dict(size=15.5, color=INK))],
], gap=12, glyph_color=RED)
card(s, 8.9, 2.55, 3.6, 2.0, "0.7866", "macro AUROC — la baseline onesta\nda cui riparte tutto", RED, num_size=40)

# ---------------------------------------------------------------- 9 · risultati GAN
s = content_slide("Fase 1 · Risultati", "La famiglia GAN: 0.8276 → 0.8378, poi il plateau", accent=RED)
pic_fit(s, FIG / "fig4_gan_per_categoria.png", 0.6, 1.6, 12.13, 4.8)
takeaway(s, [("Punte eccellenti (screw 1.000, wood 1.000, hazelnut 0.991) ma cable 0.58 e metal_nut 0.62 non si spostano. ",
              dict(size=13.5, color=INK)),
             ("E i segnali più affidabili sono proprio quelli che non si addestrano…",
              dict(size=13.5, color=DARK, bold=True))])

# ---------------------------------------------------------------- 10 · svolta patchcore
s = content_slide("Fase 2 · La svolta", "PatchCore: il modello che non si addestra")
tf = textbox(s, 0.6, 1.85, 7.4, 4.3)
bullets(tf, [
    [("Feature ImageNet congelate ", dict(size=16, color=INK, bold=True)),
     ("(WideResNet50-2) + memory bank di patch normali + distanza dal vicino più prossimo",
      dict(size=16, color=INK))],
    [("Un difetto = una patch che non assomiglia a nessuna patch normale mai vista",
      dict(size=16, color=INK, bold=True))],
    "Niente training, niente loss, niente seed: il “modello” è la memoria",
    "Costruzione del bank in secondi per categoria",
], gap=14)
card(s, 8.35, 1.95, 4.35, 1.9, "0.9051", "macro AUROC al primo tentativo:\ngià sopra qualunque GAN", BLUE, num_size=36)
p = rect(s, 8.35, 4.15, 4.35, 1.55, LIGHT)
tf = p.text_frame; tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.2)
para(tf, [("ma con buchi vistosi:\n", dict(size=12.5, color=GREY))],
     first=True, align=PP_ALIGN.CENTER, space_after=2)
para(tf, [("pill 0.61 · zipper 0.72 · capsule 0.77", dict(size=15, color=RED, bold=True))],
     align=PP_ALIGN.CENTER, space_after=0)

# ---------------------------------------------------------------- 11 · tre ingredienti
s = content_slide("Fase 2 · Tre ingredienti", "v1 → v2 → v3:  0.9051 → 0.9397 → 0.9828")
pic_fit(s, FIG / "fig7_patchcore_v1v2v3.png", 0.6, 1.6, 12.13, 4.8)
takeaway(s, [("① top-k reweighted (k=9)   ② multi-scala layer2+layer3   ③ bank pieno 70k senza coreset.  ",
              dict(size=13.5, color=DARK, bold=True)),
             ("pill: 0.534 → 0.872 cambiando solo come si interroga lo stesso bank.",
              dict(size=13.5, color=INK))])

# ---------------------------------------------------------------- 12 · paradosso coreset
s = content_slide("Fase 2 · Il paradosso", "Il bank pieno è più accurato E ~6× più veloce")
pic_fit(s, FIG / "fig8_tempi.png", 0.6, 1.6, 12.13, 4.8)
takeaway(s, [("Il tempo dominante non era la ricerca kNN: era la selezione del coreset. ",
              dict(size=13.5, color=INK)),
             ("6.3 s/categoria ", dict(size=13.5, color=GREEN, bold=True)),
             ("(hazelnut 799 s: l'unica oltre le 70k patch, il coreset scatta solo lì).",
              dict(size=13.5, color=INK))])

# ---------------------------------------------------------------- 13 · production
s = content_slide("Il sistema finale", "Production: 0.9846 macro AUROC", accent=GREEN)
cw, gap = 3.85, 0.29
card(s, 0.6, 1.9, cw, 1.6, "4 × 1.0000", "categorie perfette: bottle,\nhazelnut, leather, tile", GREEN, num_size=26)
card(s, 0.6 + cw + gap, 1.9, cw, 1.6, "11 ≥ 0.96", "delle 15 categorie", GREEN, num_size=26)
card(s, 0.6 + 2 * (cw + gap), 1.9, cw, 1.6, "0.9419", "la peggiore (screw)", BLUE, num_size=26)
tf = textbox(s, 0.6, 3.9, 12.1, 2.5)
bullets(tf, [
    [("Configurazione per-categoria: ", dict(size=15.5, color=INK, bold=True)),
     ("layer1+2+3 solo per screw (+2.7 punti); layer2+3 per tutte le altre",
      dict(size=15.5, color=INK))],
    [("Soglia calibrata: ", dict(size=15.5, color=INK, bold=True)),
     ("99° percentile su val_normal tenuto fuori dal bank (le immagini del bank hanno "
      "distanza ≈ 0 per costruzione: calibrarci sopra è una trappola)", dict(size=15.5, color=INK))],
    [("Operatività: ", dict(size=15.5, color=INK, bold=True)),
     ("build del bank in 4–8 s per categoria, inferenza deterministica, zero training",
      dict(size=15.5, color=INK))],
], gap=11, glyph_color=GREEN)

# ---------------------------------------------------------------- 14 · complementarità
s = content_slide("Confronto", "Due paradigmi complementari")
pic_fit(s, FIG / "fig6_heatmap.png", 0.6, 1.6, 12.13, 4.8)
takeaway(s, [("screw: GAN 1.000 vs production 0.942 — cable: production 0.996 vs GAN 0.58. ",
              dict(size=13.5, color=DARK, bold=True)),
             ("Misurano cose diverse: “che cosa so rigenerare” vs “che cosa ho già visto”.",
              dict(size=13.5, color=INK))])

# ---------------------------------------------------------------- 15 · evoluzione
s = content_slide("La storia in un grafico", "Da 0.7866 a 0.9846", accent=GREEN)
pic_fit(s, FIG / "fig1_evoluzione.png", 0.6, 1.6, 12.13, 4.8)
takeaway(s, [("GAN modernizzato: +5 punti di tuning accurato. Feature congelate: +16 punti in tre mosse. ",
              dict(size=13.5, color=INK)),
             ("La lezione era nei segnali ausiliari fin dall'inizio.", dict(size=13.5, color=DARK, bold=True))])

# ---------------------------------------------------------------- 16 · webapp
s = content_slide("Demo", "La webapp: i modelli, dal vivo")
pic_fit(s, FIG / "webapp_arena_gan.png", 0.6, 1.6, 12.13, 4.8)
takeaway(s, [("Test Arena con 5 varianti live ", dict(size=13.5, color=DARK, bold=True)),
             ("(production, PatchCore v1/v2, OCGAN final/optv2 dai checkpoint originali) — "
              "streaming SSE, Evaluation Lab con questi stessi numeri, Dataset Explorer.",
              dict(size=13.5, color=INK))])

# ---------------------------------------------------------------- 17 · conclusioni
s = content_slide("Conclusioni", "Cinque lezioni dal percorso")
tf = textbox(s, 0.6, 1.8, 12.1, 4.5)
lessons = [
    ("La fusione batte il singolo segnale", " — +0.22 AUROC da una regressione logistica, non da un'architettura"),
    ("Le feature pre-addestrate sono il moltiplicatore", " — portate all'estremo, sono PatchCore"),
    ("Il rigore è una feature del risultato", " — flag morti, score dannosi e trappole di calibrazione si trovano solo se i numeri devono essere onesti"),
    ("I log dei tempi sono diagnostica", " — il coreset inutile e l'outlier hazelnut sono usciti dai secondi, non dall'AUROC"),
    ("Le idee sopravvivono ai modelli", " — di OCGAN resta l'atteggiamento: cercare attivamente dove il modello sbaglia"),
]
for i, (head, rest) in enumerate(lessons):
    para(tf, [(f"{i+1}.  ", dict(size=17, color=GREEN, bold=True)),
              (head, dict(size=17, color=DARK, bold=True)),
              (rest, dict(size=15.5, color=INK))],
         first=(i == 0), space_after=14)
takeaway(s, [("Grazie!  ", dict(size=14, color=DARK, bold=True)),
             ("Tutti i dettagli, le tabelle per-categoria e le ablation complete: RELAZIONE.pdf (39 pagine).",
              dict(size=13.5, color=INK))], y=6.55)

# ---------------------------------------------------------------- 18 · deployment, verifica, ricalibrazione
s = content_slide("Deployment · Verifica & ricalibrazione", "Verifica live, soglia per categoria, localizzazione pixel-level")
tf = textbox(s, 0.6, 1.7, 12.13, 3.4)
bullets(tf, [
    [("Verifica end-to-end: ", dict(size=15, color=INK, bold=True)),
     ("verify_all.py riesegue dal vivo ogni variante x categoria — drift nullo sulle 45 combinazioni", dict(size=15, color=INK))],
    [("Bug della soglia su screw: ", dict(size=15, color=INK, bold=True)),
     ("AUROC 0.93 ma accuratezza arena 0.44 — la soglia p99 dei normali e' troppo alta per una categoria rotation-variant", dict(size=15, color=INK))],
    [("Soglia best-F1 per categoria: ", dict(size=15, color=INK, bold=True)),
     ("override non distruttivo: screw 0.44->0.92, capsule 0.61->0.98, zero regressioni; AUROC invariata + slider interattivo in arena", dict(size=15, color=INK))],
    [("Localizzazione pixel-level: ", dict(size=15, color=INK, bold=True)),
     ("anomaly-map vs maschere ground-truth: macro pixel-AUROC 0.9714 e AUPRO@30% 0.9127 (metrica ufficiale MVTec), accanto allo 0.9846 image-level", dict(size=15, color=INK))],
    [("Calibrazione onesta (soglia + probabilita'): ", dict(size=15, color=INK, bold=True)),
     ("best-F1 e' un oracle -> accuratezza held-out 0.954 vs 0.968; lo score grezzo non e' una probabilita' (ECE 0.32) -> calibrato Platt/isotonica ECE 0.03 (held-out)", dict(size=15, color=INK))],
    [("TTA rotazioni, risultato negativo: ", dict(size=15, color=INK, bold=True)),
     ("peggiora screw (0.93->0.89), il bank gia' codifica le orientazioni. Non integrata", dict(size=15, color=INK))],
], gap=8)
cw, gap = 3.4, 0.45
x0 = (SW - 3 * cw - 2 * gap) / 2
card(s, x0, 5.55, cw, 1.0, "0.9714", "macro pixel-AUROC", GREEN)
card(s, x0 + cw + gap, 5.55, cw, 1.0, "0.9127", "AUPRO@30% (ufficiale MVTec)", GREEN)
card(s, x0 + 2 * (cw + gap), 5.55, cw, 1.0, "0.954", "accuratezza onesta (held-out)", BLUE)
takeaway(s, [("Misurare l'operating point, non solo il ranking: l'AUROC era giusta, la soglia no.", dict(size=14, color=DARK, bold=True))])

# ---------------------------------------------------------------- 19 · scelte di efficienza
s = content_slide("Scelte di efficienza", "Cosa NON abbiamo aggiunto — e perché", accent=GREEN)
tf = textbox(s, 0.6, 1.75, 12.13, 3.0)
bullets(tf, [
    [("Niente FAISS / ANN: ", dict(size=15.5, color=INK, bold=True)),
     ("a questa dimensione del bank la ricerca esatta del vicino costa secondi per categoria; "
      "un indice approssimato aggiungerebbe una dipendenza e un rischio di accuratezza per zero guadagno", dict(size=15.5, color=INK))],
    [("Niente backbone piu' pesante: ", dict(size=15.5, color=INK, bold=True)),
     ("WideResNet50-2 congelata raggiunge gia' 0.9846 macro; l'intero sistema addestra e serve su una GPU da 4 GB", dict(size=15.5, color=INK))],
    [("Ensemble GAN + PatchCore, misurato non assunto: ", dict(size=15.5, color=INK, bold=True)),
     ("peso di fusione scelto su fold di calibrazione e validato su fold held-out — nessun guadagno materiale "
      "sul macro AUROC: il PatchCore da solo domina, quindi spediamo il sistema piu' semplice", dict(size=15.5, color=INK))],
], gap=12, glyph_color=GREEN)
cw, gap = 3.85, 0.29
card(s, 0.6, 5.15, cw, 1.25, "4 GB", "tutta l'inferenza su Quadro T1000", GREEN, num_size=30)
card(s, 0.6 + cw + gap, 5.15, cw, 1.25, "~6 s", "build del bank per categoria", GREEN, num_size=30)
card(s, 0.6 + 2 * (cw + gap), 5.15, cw, 1.25, "0.9846", "macro AUROC, gia' al target", BLUE, num_size=30)
takeaway(s, [("Efficienza come scelta progettuale: gli obiettivi alti sono raggiunti senza complessita' aggiuntiva.",
              dict(size=14, color=DARK, bold=True))])

prs.save(OUT)
print(f"ok: {OUT} ({OUT.stat().st_size} bytes, {len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst)} slide)")
