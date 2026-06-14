# -*- coding: utf-8 -*-
"""Verifica strutturale di PRESENTAZIONE.pptx: slide, immagini, testi."""
import sys
from pathlib import Path

from pptx import Presentation

sys.stdout.reconfigure(encoding="utf-8", errors="replace")

PPTX = Path(__file__).resolve().parent / "PRESENTAZIONE.pptx"

p = Presentation(PPTX)
print(f"file: {PPTX.name} ({PPTX.stat().st_size} bytes)")
print(f"slide totali: {len(p.slides)}")
tot_pics = 0
for i, s in enumerate(p.slides, 1):
    pics = [sh for sh in s.shapes if sh.shape_type == 13]
    tot_pics += len(pics)
    texts = []
    for sh in s.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip().replace("\n", " | ")
            if t:
                texts.append(t)
    title = texts[1] if len(texts) > 1 else (texts[0] if texts else "(vuota)")
    print(f"  {i:2d}. img={len(pics)}  txt={len(texts):2d}  {title[:72]}")
print(f"immagini totali: {tot_pics}")
